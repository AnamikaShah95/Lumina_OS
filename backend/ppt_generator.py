import json
import os
import time
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from config.settings import GEMINI_API_KEY
from backend.presentation_engine import PresentationPayload

class LuminaPPTGenerator:
    def __init__(self):
        if not GEMINI_API_KEY:
            raise ValueError("❌ Critical Error: GEMINI_API_KEY is missing in environment.")
        self.client = genai.Client(api_key=GEMINI_API_KEY)
        self.model_name = 'gemini-2.5-flash'

        # Premium Corporate Design Tokens Palette
        self.COLOR_BG = RGBColor(20, 21, 30)           # Shadow Eclipse Navy
        self.COLOR_ACCENT_BG = RGBColor(28, 30, 43)    # Subtle Accent Backdrop Mesh
        self.COLOR_CARD_BG = RGBColor(35, 38, 55)      # Glassmorphism Card Fill
        self.COLOR_TEXT_MAIN = RGBColor(248, 248, 242) # Pearl White Core
        self.COLOR_ACCENT = RGBColor(189, 147, 249)    # Neon Kawaii Purple
        self.COLOR_MUTED = RGBColor(139, 233, 253)     # Cyber Cyan Core Accent

    def transform_summary_to_slides(self, video_title: str, summary_text: str, target_slides: int, audience_level: str) -> dict:
        print(f"📊 [PPT Engine]: Transforming video parameters inside multi-layer structured pipelines...")
        prompt = f"""
        You are an expert Enterprise AI Slide Architect. Parse the input summaries and structure it into a clean JSON schema payload matching the exact shape.

        CRITICAL CONSTRAINTS:
        1. Create exactly {target_slides} individual items inside the 'slides' list block array.
        2. Set the 'total_slides' field integer value to exactly equal {target_slides}.
        3. Target Depth Profile: '{audience_level}'.
        4. Inside every slide object, the 'bullet_points' array MUST contain between 3 to 5 clear strings.

        Metadata Context:
        Master Title: {video_title}
        Source Text Component Payload:
        {summary_text}
        """
        try:
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=prompt,
                config={
                    'response_mime_type': 'application/json',
                    'response_schema': PresentationPayload,
                }
            )
            structured_data = json.loads(response.text)
            print(f"✅ [PPT Engine]: LLM Structural Payload Arrays locked successfully.")
            return {"status": "success", "data": structured_data, "error": None}
        except Exception as e:
            return {"status": "failed", "data": None, "error": str(e)}

    def _apply_background_with_geometric_accents(self, slide):
        base_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
        base_shape.fill.solid()
        base_shape.fill.fore_color.rgb = self.COLOR_BG
        base_shape.line.fill.background()

        accent_block = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(5.5), Inches(0), Inches(4.5), Inches(4.0))
        accent_block.fill.solid()
        accent_block.fill.fore_color.rgb = self.COLOR_ACCENT_BG
        accent_block.line.fill.background()
        accent_block.rotation = 180 

        # Day 25 Structural Move: Removed old static line divider from here 
        # to prevent lines cutting down middle of text wraps during multi-line headers.

    def _draw_content_card(self, slide, left, top, width, height):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLOR_CARD_BG
        card.line.fill.background()
        
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.COLOR_MUTED
        accent_bar.line.fill.background()

    def _compute_adaptive_font_size(self, text_string: str) -> int:
        char_count = len(text_string)
        if char_count <= 40: return 38
        elif char_count <= 75: return 26
        else: return 20

    def _generate_mock_visual_placeholder(self, slide, left, top, width, height, placeholder_text: str):
        img_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = RGBColor(43, 47, 68)
        img_box.line.color.rgb = self.COLOR_MUTED
        img_box.line.width = Pt(1)
        
        tx_box = slide.shapes.add_textbox(left, top + (height / 2) - Inches(0.4), width, Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"⚙️ [ VISUAL NODE CONTEXT ]\n{placeholder_text.upper()}"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(10)
        p.font.color.rgb = self.COLOR_MUTED
        p.alignment = PP_ALIGN.CENTER

    def _inject_fade_transition(self, slide):
        slide_element = slide.element
        transition_node = OxmlElement('p:transition')
        transition_node.set('spd', 'fast')
        fade_node = OxmlElement('p:fade')
        transition_node.append(fade_node)
        slide_element.append(transition_node)

    def _format_rich_bullet_point(self, paragraph, text: str, font_size_pt: int = 13):
        paragraph.text = "" 
        paragraph.line_spacing = 1.3               
        paragraph.space_before = Pt(4)             
        paragraph.space_after = Pt(12)             
        
        marker_run = paragraph.add_run()
        marker_run.text = "▪  "
        marker_run.font.name = 'Segoe UI'
        marker_run.font.size = Pt(font_size_pt - 1)
        marker_run.font.bold = True
        marker_run.font.color.rgb = self.COLOR_MUTED 
        
        body_run = paragraph.add_run()
        body_run.text = text
        body_run.font.name = 'Calibri'
        body_run.font.size = Pt(font_size_pt)
        body_run.font.color.rgb = self.COLOR_TEXT_MAIN

    def _inject_dynamic_progress_tracker(self, slide, current_idx: int, total_slides: int):
        if total_slides <= 0: return
        base_track_width = 8.5
        calculated_progress_ratio = current_idx / total_slides
        progress_bar_width = Inches(base_track_width * calculated_progress_ratio)
        
        progress_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(7.25), progress_bar_width, Inches(0.04))
        progress_shape.fill.solid()
        progress_shape.fill.fore_color.rgb = self.COLOR_ACCENT 
        progress_shape.line.fill.background()
        
        counter_box = slide.shapes.add_textbox(Inches(8.25), Inches(6.8), Inches(1.0), Inches(0.4))
        p = counter_box.text_frame.paragraphs[0]
        p.text = f"{str(current_idx).zfill(2)} / {str(total_slides).zfill(2)}"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(11)
        p.font.color.rgb = self.COLOR_MUTED 
        p.alignment = PP_ALIGN.RIGHT

    def generate_actual_pptx(self, structured_data: dict, output_filename: str = "presentation.pptx") -> str:
        print("🛠️ [Safety Wrap Architect]: Compiling context-aware auto-wrapping header bounds...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # ----------------------------------------------------------------------
        # TITLE SLIDE ARCHITECTURE
        # ----------------------------------------------------------------------
        slide = prs.slides.add_slide(blank_layout)
        self._apply_background_with_geometric_accents(slide)
        self._inject_fade_transition(slide)

        master_topic = structured_data.get("topic", "Lumina OS Visual Core").upper()
        dynamic_title_font_size = self._compute_adaptive_font_size(master_topic)

        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(8.5), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = master_topic
        p.font.name = 'Segoe UI'
        p.font.size = Pt(dynamic_title_font_size)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_ACCENT

        p2 = tf.add_paragraph()
        p2.text = f"Automated Header Bounding Box Optimization Pipeline\nDynamic Spatial Offset Matrix Layers | Nodes Count: {structured_data.get('total_slides', 0)}"
        p2.font.name = 'Segoe UI'
        p2.font.size = Pt(13)
        p2.font.color.rgb = self.COLOR_MUTED
        p2.space_before = Pt(20)

        # ----------------------------------------------------------------------
        # THREE-ZONE SLIDES PACK WITH SAFETY HEIGHT ADAPTIVE LAYOUTS
        # ----------------------------------------------------------------------
        slides_list = structured_data.get("slides", [])
        total_slides_count = int(structured_data.get("total_slides", len(slides_list)))

        for idx, slide_data in enumerate(slides_list, start=1):
            slide = prs.slides.add_slide(blank_layout)
            self._apply_background_with_geometric_accents(slide)
            self._inject_fade_transition(slide)
            self._inject_dynamic_progress_tracker(slide, current_idx=idx, total_slides=total_slides_count)

            slide_title_text = f"{slide_data.get('slide_number')}. {slide_data.get('title')}"
            header_chars_count = len(slide_title_text)

            # --- Day 25 Mathematical Spatial Optimization Calculations ---
            # Automatically adjusting height bounds and layout shifts matching title string volumes
            if header_chars_count <= 45:
                # Standard short header title configuration bounds
                title_allocated_height = Inches(0.8)
                content_top_offset = Inches(1.6)
                card_rendering_height = Inches(5.0)
                font_pt_size = 24
            else:
                # Complex multi-line long header wrap state configuration bounds
                title_allocated_height = Inches(1.4)
                content_top_offset = Inches(2.2)       # Pushing the blocks layout down by 0.6 inches
                card_rendering_height = Inches(4.4)    # Compacting card height limits safely to save padding
                font_pt_size = 20                       # Slicing scale dimensions to balance wrap tracks

            # Render the wrap-enabled adaptive title box
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8.5), title_allocated_height)
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_title_text
            p_title.font.name = 'Segoe UI'
            p_title.font.size = Pt(font_pt_size)
            p_title.font.bold = True
            p_title.font.color.rgb = self.COLOR_ACCENT

            # Adaptive Dynamic Tech Divider Track Line (Injected dynamically directly below wrap boundaries)
            line_top_position = Inches(0.4) + title_allocated_height + Inches(0.1)
            divider_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), line_top_position, Inches(8.5), Inches(0.02))
            divider_line.fill.solid()
            divider_line.fill.fore_color.rgb = RGBColor(45, 49, 69)
            divider_line.line.fill.background()

            # Core column track geometry distributions
            points = slide_data.get("bullet_points", [])
            mid_index = (len(points) + 1) // 2
            left_points = points[:mid_index]
            right_points = points[mid_index:]

            col_width = Inches(2.7)
            img_width = Inches(2.6)

            # ZONE 1: Context Graphic Placeholder (Height and top offset adapts dynamically)
            img_left_start = Inches(0.75)
            self._generate_mock_visual_placeholder(
                slide=slide, left=img_left_start, top=content_top_offset, width=img_width, height=card_rendering_height, 
                placeholder_text=slide_data.get('title', 'Analytics Data Core')
            )

            # ZONE 2: TEXT CONTENT COLUMN A (Middle Track)
            col_left_start = Inches(3.6)
            self._draw_content_card(slide, col_left_start, content_top_offset, col_width, card_rendering_height)
            
            body_box_left = slide.shapes.add_textbox(col_left_start + Inches(0.15), content_top_offset + Inches(0.15), col_width - Inches(0.2), card_rendering_height - Inches(0.3))
            tf_left = body_box_left.text_frame
            tf_left.word_wrap = True

            for i, pt_text in enumerate(left_points):
                p_pt = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
                self._format_rich_bullet_point(p_pt, pt_text, font_size_pt=13)

            # ZONE 3: TEXT CONTENT COLUMN B (Right Margin Track)
            col_right_start = Inches(6.55)
            self._draw_content_card(slide, col_right_start, content_top_offset, col_width, card_rendering_height)
            
            body_box_right = slide.shapes.add_textbox(col_right_start + Inches(0.15), content_top_offset + Inches(0.15), col_width - Inches(0.2), card_rendering_height - Inches(0.3))
            tf_right = body_box_right.text_frame
            tf_right.word_wrap = True

            for i, pt_text in enumerate(right_points):
                p_pt = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
                self._format_rich_bullet_point(p_pt, pt_text, font_size_pt=13)

        output_dir = "output"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        final_path = os.path.join(output_dir, output_filename)
        prs.save(final_path)
        print(f"💾 [Safety Wrap Architect]: Adaptive header optimized presentations compiled cleanly at: {final_path}")
        return final_path